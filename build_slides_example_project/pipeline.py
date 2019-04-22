import os
import luigi
import pptx
from typing import List
from luigi.contrib.external_program import ExternalProgramTask
from datetime import datetime
from dateutil import tz

class PrintDate(luigi.Task):
    pptx_filename: str = luigi.Parameter()
    workdir: str = luigi.Parameter()

    def requires(self):
        return None

    def output(self):
        filename = self.pptx_filename.replace('.pptx', '_wdate.pptx')
        return luigi.LocalTarget(os.path.join(self.workdir, filename))

    def run(self):
        from_zone = tz.gettz('UTC')
        to_zone = tz.tzlocal()

        # Convert utc time to local
        build_time = datetime.utcnow()
        build_time = build_time.replace(tzinfo=from_zone)
        build_time = build_time.astimezone(to_zone)

        build_timestamp: str = build_time.strftime('%d, %b %Y %H:%M:%S %p')

        # Replace the placeholder with the timestamp
        slide: int = int(self.pptx_filename.split('.pptx')[0].split('_')[-1])
        pst: pptx.Presentation = pptx.Presentation(os.path.join(self.workdir, self.pptx_filename))

        first_slide: pptx.slide.Slide = pst.slides[slide]
        shapes: List[pptx.shapes.base] = first_slide.shapes
        paragraphs = [shape.text_frame for shape in shapes if shape.has_text_frame]

        for paragraph in paragraphs:
            paragraph.text = paragraph.text.replace('[date]', build_timestamp)
            paragraph.text = paragraph.text.replace('[title]', 'Pipelines with Luigi')
            paragraph.text = paragraph.text.replace('[author]', 'Alejandro Rodríguez Díaz')

        pst.save(self.output().path)

class ExtraProcessing(luigi.Task):
    pptx_filename: str = luigi.Parameter()
    workdir: str = luigi.Parameter()

    def requires(self):
        return PrintDate(self.pptx_filename, self.workdir)

    def output(self):
        filename = self.pptx_filename.replace('.pptx', '_processed.pptx')
        return luigi.LocalTarget(os.path.join(self.workdir, filename))

    def run(self):
        pst: pptx.Presentation = pptx.Presentation(self.input().path)
        pst.save(self.output().path)


class Pptx2Pdf(ExternalProgramTask):
    pptx_filename: str = luigi.Parameter()
    workdir: str = luigi.Parameter()

    def requires(self):
        return ExtraProcessing(self.pptx_filename, self.workdir)

    def output(self):
        filename = self.pptx_filename.replace('.pptx', '_processed.pptx')
        pdf_filename = filename.replace('.pptx', '.pdf')
        return luigi.LocalTarget(os.path.join(self.workdir, pdf_filename))

    def program_args(self):
        filename = self.pptx_filename.replace('.pptx', '_processed.pptx')
        pdf_filename = filename.replace('.pptx', '.pdf')

        return [
            "docker",
            "run",
            "--rm",
            "-v",
            f"{self.workdir}:/data",
            "seguins/soffice",
            "bash", 
            "-c" ,
            "soffice --headless --convert-to pdf:impress_pdf_Export /data/" + \
                f"{filename} && cp {pdf_filename} /data"
        ]

class MergeSlides(ExternalProgramTask):
    pptx_filename: str = luigi.Parameter()
    workdir: str = luigi.Parameter()

    def requires(self):
        target_from_index = lambda i: self.pptx_filename.replace('.pptx', f'_raw_{i}.pptx')

        pst: pptx.Presentation = pptx.Presentation(os.path.join(self.workdir, self.pptx_filename))
        for i in range(len(pst.slides)):
            yield Pptx2Pdf(workdir=self.workdir, pptx_filename=target_from_index(i))

    def output(self):
        filename = self.pptx_filename.replace('.pptx', f'.pdf')
        return luigi.LocalTarget(os.path.join(self.workdir, filename))

    def program_args(self):
        slides: List[str] = list(f.path for f in self.input())
        args: List[str] = ['pdfunite']
        args.extend(slides)
        args.append(self.output().path)

        return args

class ExtractSlides(luigi.Task):
    pptx_filename: str = luigi.Parameter()
    workdir: str = luigi.Parameter()

    def requires(self):
        return None

    def output(self):
        target_from_index = lambda i: luigi.LocalTarget( \
            os.path.join( \
                self.workdir,  
                self.pptx_filename.replace('.pptx', f'_raw_{i}.pptx')))

        pst: pptx.Presentation = pptx.Presentation(os.path.join(self.workdir, self.pptx_filename))
        return {i:target_from_index(i) for i in range(len(pst.slides))}

    def run(self):
        pst: pptx.Presentation = pptx.Presentation(os.path.join(self.workdir, self.pptx_filename))

        for slide in pst.slides:
            slide._element.set('show', '0')

        for i in range(len(pst.slides)):
            pst.slides[i]._element.set('show', '1')
            filename = self.pptx_filename.replace('.pptx', f'_raw_{i}.pptx')
            pst.save(os.path.join(self.workdir, filename))
            pst.slides[i]._element.set('show', '0')

class Pipeline(luigi.Task):
    pptx_filename: str = luigi.Parameter()
    workdir: str = luigi.Parameter()

    def requires(self):
        return ExtractSlides(workdir=self.workdir, pptx_filename=self.pptx_filename)

    def output(self):
        filename = self.pptx_filename.replace('.pptx', f'.pdf')
        return luigi.LocalTarget(os.path.join(self.workdir, filename))

    def run(self):
        yield MergeSlides(workdir=self.workdir, pptx_filename=self.pptx_filename)



if __name__ == '__main__':
    luigi.build([Pipeline(workdir=os.path.abspath('./slides'), pptx_filename='base.pptx')], workers=6)
